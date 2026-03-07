# Necesitas instalar: pip install python-pptx
from generar_presentacion import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Diapositiva 1
title_slide = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide)
slide.shapes.title.text = "Proyecto Eficiencia Energética"
slide.placeholders[1].text = "Carmen Canales - SBD Víctor\nAutomatización con Dataflows"

# Diapositiva 2
add_slide("Arquitectura de Datos", 
          "• Origen: SharePoint Online.\n• Motor: Power BI Dataflows (Procesamiento en la nube).\n• Conector: SharePoint.Contents.")

# Diapositiva 3
add_slide("Transformación M (Power Query)", 
          "• Sincronización horaria (0-23) para PVPC.\n• Consolidación automática de múltiples CSVs.\n• Generación de claves subrogadas (DateKey).")

# Diapositiva 4
add_slide("Modelo de Datos", 
          "• Esquema en Estrella.\n• Dimensión: Calendario (Fechas).\n• Hechos: Consumo, Precios y Generación Enphase.")

# Diapositiva 5
add_slide("Visualización y DAX", 
          "• Comparativa Escenario 1 vs Escenario 2.\n• Cálculo de ahorro horario.\n• Interactividad total mediante segmentadores.")

prs.save('Proyecto_Carmen_Eficiencia.pptx')
print("¡PowerPoint creado con éxito!")